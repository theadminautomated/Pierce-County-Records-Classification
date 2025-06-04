#!/usr/bin/env python3
"""
File Scanner Module for Electronic Records Classification
---------------------------------------------------------
Efficiently discovers and categorizes files for processing with proper architecture.
"""

import os
import sys
import subprocess
from pathlib import Path
from typing import Dict, Set, List, Iterator, Tuple, Union
from dataclasses import dataclass
import datetime
import logging

try:
    from docx import Document
except Exception:  # pragma: no cover - optional dependency
    Document = None

try:
    import pdfplumber
except Exception:  # pragma: no cover - optional dependency
    pdfplumber = None

try:
    import PyPDF2
except Exception:  # pragma: no cover - optional dependency
    PyPDF2 = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency
    Presentation = None

try:
    import openpyxl
except Exception:  # pragma: no cover - optional dependency
    openpyxl = None

try:
    import xlrd
except Exception:  # pragma: no cover - optional dependency
    xlrd = None

try:
    from PIL import Image
    import pytesseract
except Exception:  # pragma: no cover - optional dependency
    Image = None
    pytesseract = None

from pathlib import Path
from typing import List, Dict, Set, Optional, Iterator, Tuple, Union
from dataclasses import dataclass
import datetime
import logging
import subprocess
 main

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# File type constants - matches main classification module
INCLUDE_EXT: Set[str] = frozenset({
    '.txt', '.csv', '.docx', '.xlsx', '.pptx', '.pdf', '.html', '.htm', '.md',
    '.rtf', '.odt', '.xml', '.json', '.yaml', '.yml', '.log', '.tsv'
})

EXCLUDE_EXT: Set[str] = frozenset({
    '.tmp', '.bak', '.old', '.zip', '.rar', '.tar', '.gz', '.7z',
    '.exe', '.dll', '.sys', '.iso', '.dmg', '.apk', '.msi', '.ps1', '.psd1',
    '.psm1', '.db', '.mdb', '.accdb'
})

@dataclass
class FileInfo:
    """Information about a discovered file."""
    path: Path
    size_bytes: int
    modified_time: datetime.datetime
    extension: str
    category: str  # 'destroy', 'analyze', 'skip'
    reason: str    # Why it was categorized this way

class FileScanner:
    """
    Handles file discovery and categorization for classification.
    """
    
    def __init__(self, include_ext: Set[str] = None, exclude_ext: Set[str] = None):
        self.include_ext = include_ext or INCLUDE_EXT
        self.exclude_ext = exclude_ext or EXCLUDE_EXT
        self.destroy_threshold = datetime.datetime.now() - datetime.timedelta(days=6 * 365)
    
    def scan_directory(self, directory_path: Union[str, Path]) -> Iterator[FileInfo]:
        """
        Scan a directory and yield FileInfo objects for all discovered files.
        
        Args:
            directory_path: Path to directory to scan
            
        Yields:
            FileInfo objects for each discovered file
        """
        directory = Path(directory_path)
        
        if not directory.exists():
            raise ValueError(f"Directory does not exist: {directory}")
        
        if not directory.is_dir():
            raise ValueError(f"Path is not a directory: {directory}")
        
        logger.info(f"Scanning directory: {directory}")
        
        # Use rglob to recursively find all files
        for file_path in directory.rglob('*'):
            if not file_path.is_file():
                continue
                
            # Skip hidden files and temp files
            if file_path.name.startswith('.') or file_path.name.startswith('~$'):
                continue
            
            try:
                file_info = self._analyze_file(file_path)
                yield file_info
            except Exception as e:
                logger.warning(f"Error analyzing file {file_path}: {e}")
                # Still yield the file with error info
                yield FileInfo(
                    path=file_path,
                    size_bytes=0,
                    modified_time=datetime.datetime.now(),
                    extension=file_path.suffix.lower(),
                    category='skip',
                    reason=f"Error analyzing file: {e}"
                )
    
    def _analyze_file(self, file_path: Path) -> FileInfo:
        """
        Analyze a single file and determine its category.
        
        Args:
            file_path: Path to the file
            
        Returns:
            FileInfo object with categorization
        """
        try:
            stat_info = file_path.stat()
            modified_time = datetime.datetime.fromtimestamp(stat_info.st_mtime)
            extension = file_path.suffix.lower()
            
            # Determine category
            category, reason = self._categorize_file(file_path, modified_time, extension)
            
            return FileInfo(
                path=file_path,
                size_bytes=stat_info.st_size,
                modified_time=modified_time,
                extension=extension,
                category=category,
                reason=reason
            )
            
        except Exception as e:
            logger.error(f"Error getting file stats for {file_path}: {e}")
            raise
    
    def _categorize_file(self, file_path: Path, modified_time: datetime.datetime, extension: str) -> Tuple[str, str]:
        """
        Categorize a file based on age, type, and other criteria.
        
        Args:
            file_path: Path to the file
            modified_time: When the file was last modified
            extension: File extension
            
        Returns:
            Tuple of (category, reason)
        """
        # Check if file is old enough for automatic destroy
        if modified_time < self.destroy_threshold:
            return 'destroy', 'Older than 6 years - automatic destroy'
        
        # Check if file type is excluded
        if extension in self.exclude_ext:
            return 'skip', f'Excluded file type: {extension}'
        
        # Check if file type is not in include list
        if extension not in self.include_ext:
            return 'skip', f'Unsupported file type: {extension}'
        
        # File needs analysis
        return 'analyze', 'Supported file type within retention period'
    
    def get_file_counts(self, directory_path: Union[str, Path]) -> Dict[str, int]:
        """
        Get counts of files by category without yielding individual files.
        
        Args:
            directory_path: Path to directory to scan
            
        Returns:
            Dictionary with counts for each category
        """
        counts = {'destroy': 0, 'analyze': 0, 'skip': 0, 'total': 0}
        
        for file_info in self.scan_directory(directory_path):
            counts[file_info.category] += 1
            counts['total'] += 1
        
        return counts
    
    def filter_files_by_category(
        self, 
        directory_path: Union[str, Path], 
        category: str
    ) -> Iterator[FileInfo]:
        """
        Get only files that match a specific category.
        
        Args:
            directory_path: Path to directory to scan
            category: Category to filter by ('destroy', 'analyze', 'skip')
            
        Yields:
            FileInfo objects matching the category
        """
        for file_info in self.scan_directory(directory_path):
            if file_info.category == category:
                yield file_info

# Convenience functions for backward compatibility
def scan_files(folder: Union[str, Path], include_ext: Set[str] = None) -> Iterator[Path]:
    """
    Simple file scanner that yields Path objects.
    
    Args:
        folder: Directory to scan
        include_ext: Set of extensions to include
        
    Yields:
        Path objects for discovered files
    """
    scanner = FileScanner(include_ext=include_ext)
    for file_info in scanner.scan_directory(folder):
        if file_info.category == 'analyze':
            yield file_info.path

def get_file_categories(folder: Union[str, Path]) -> Dict[str, List[Path]]:
    """
    Get files categorized by processing type.
    
    Args:
        folder: Directory to scan
        
    Returns:
        Dictionary with lists of files by category
    """
    scanner = FileScanner()
    categories = {'destroy': [], 'analyze': [], 'skip': []}
    
    for file_info in scanner.scan_directory(folder):
        categories[file_info.category].append(file_info.path)
    
    return categories

def scan_files(folder: Union[str, Path], include_ext: Set[str] = INCLUDE_EXT) -> Iterator[Path]:
    """
    Scan a folder and yield Path objects for files with ONLY the allowed extensions.
    Ignores hidden/system files.
    """
    folder_path = Path(folder)
    for f in folder_path.rglob('*'):
        if not f.is_file():
            continue
        suffix = f.suffix.lower()
        if suffix in include_ext and not f.name.startswith('.') and not f.name.startswith('~$'):
            yield f

 def extract_file_content(f: Path, max_chars: int = 4000) -> str:
    """Extract and clean text from the given file."""
    suffix = f.suffix.lower()
    try:
        if suffix == ".txt":
            return _extract_txt(f, max_chars)
        if suffix == ".pdf":
            return _extract_pdf(f, max_chars)
        if suffix == ".docx":
            return _extract_docx(f, max_chars)
        if suffix == ".doc":
            return _extract_doc(f, max_chars)
        if suffix == ".pptx":
            return _extract_pptx(f, max_chars)
        if suffix == ".xlsx":
            return _extract_xlsx(f, max_chars)
        if suffix == ".xls":
            return _extract_xls(f, max_chars)
        if suffix in {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp"}:
            return _extract_image(f, max_chars)
        return f"[Unsupported file type: {suffix}]"
    except Exception as e:  # pragma: no cover - broad safety net
        return f"[Error extracting content: {str(e)}]"


def _extract_txt(file_path: Path, max_chars: int) -> str:
    """Return cleaned text from a plain text file."""
    for encoding in ("utf-8", "latin-1"):
        try:
            with file_path.open("r", encoding=encoding, errors="ignore") as fh:
                return _clean_text(fh.read(max_chars))
        except UnicodeDecodeError:
            continue
        except Exception as exc:  # pragma: no cover - unexpected read issues
            return f"[Error reading file: {exc}]"
    return f"[Unreadable file: {file_path.suffix}]"


def _extract_pdf(file_path: Path, max_chars: int) -> str:
    """Extract text from a PDF using available backends and OCR as fallback."""
    text = ""
    if pdfplumber:
        try:
            with pdfplumber.open(str(file_path)) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
                    if len(text) >= max_chars:
                        break
        except Exception:
            text = ""
    if not text and PyPDF2:
        try:
            with open(file_path, "rb") as fp:
                reader = PyPDF2.PdfReader(fp)
                for page in reader.pages:
                    text += page.extract_text() or ""
                    if len(text) >= max_chars:
                        break
        except Exception:
            text = ""
    if not text and pytesseract and Image:
        try:
            import pdf2image

            images = pdf2image.convert_from_path(str(file_path))
            for img in images:
                text += pytesseract.image_to_string(img)
                if len(text) >= max_chars:
                    break
        except Exception:
            text = ""
    return _clean_text(text)[:max_chars] if text else "[Could not extract text from PDF]"


def _extract_docx(file_path: Path, max_chars: int) -> str:
    """Extract text from a modern Word document."""
    if not Document:
        return "[python-docx not installed]"
    try:
        doc = Document(str(file_path))
        text = "\n".join(p.text for p in doc.paragraphs)
        return _clean_text(text)[:max_chars]
    except Exception as exc:
        return f"[Error reading DOCX: {exc}]"


def _extract_doc(file_path: Path, max_chars: int) -> str:
    """Extract text from a legacy Word document using antiword."""
    try:
        result = subprocess.run(
            ["antiword", str(file_path)], stdout=subprocess.PIPE, check=True
        )
        text = result.stdout.decode("utf-8", errors="ignore")
        return _clean_text(text)[:max_chars]
    except FileNotFoundError:
        return "[antiword not installed for .doc]"
    except subprocess.CalledProcessError as exc:
        return f"[Error reading DOC: {exc}]"


def _extract_pptx(file_path: Path, max_chars: int) -> str:
    """Extract text from a PowerPoint presentation."""
    if not Presentation:
        return "[python-pptx not installed]"
    try:
        prs = Presentation(str(file_path))
        text: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
                    if sum(len(t) for t in text) >= max_chars:
                        break
            if sum(len(t) for t in text) >= max_chars:
                break
        return _clean_text("\n".join(text))[:max_chars]
    except Exception as exc:
        return f"[Error reading PPTX: {exc}]"


def _extract_xlsx(file_path: Path, max_chars: int) -> str:
    """Extract text from an Excel workbook."""
    if not openpyxl:
        return "[openpyxl not installed]"
    try:
        wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
        text: list[str] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        text.append(str(cell))
                        if sum(len(t) for t in text) >= max_chars:
                            break
                if sum(len(t) for t in text) >= max_chars:
                    break
            if sum(len(t) for t in text) >= max_chars:
                break
        return _clean_text(" ".join(text))[:max_chars]
    except Exception as exc:
        return f"[Error reading XLSX: {exc}]"


def _extract_xls(file_path: Path, max_chars: int) -> str:
    """Extract text from a legacy Excel workbook."""
    if not xlrd:
        return "[xlrd not installed]"
    try:
        wb = xlrd.open_workbook(str(file_path))
        text: list[str] = []
        for sheet in wb.sheets():
            for row_idx in range(sheet.nrows):
                row = sheet.row_values(row_idx)
                text.append(" ".join(str(cell) for cell in row))
                if sum(len(t) for t in text) >= max_chars:
                    break
            if sum(len(t) for t in text) >= max_chars:
                break
        return _clean_text("\n".join(text))[:max_chars]
    except Exception as exc:
        return f"[Error reading XLS: {exc}]"


def _extract_image(file_path: Path, max_chars: int) -> str:
    """OCR text from an image file if Pillow and Tesseract are available."""
    if not (Image and pytesseract):
        return "[OCR dependencies not installed]"
    try:
        img = Image.open(str(file_path))
        text = pytesseract.image_to_string(img)
        return _clean_text(text)[:max_chars]
    except Exception as exc:
        return f"[Error reading image: {exc}]"

def extract_file_content(f: Path, max_chars: int = 4000) -> str:
    """
    Extract text content from a file, using OCR/parsers for binary formats.
    Returns up to max_chars of cleaned text.
    """
    suffix = f.suffix.lower()
    try:
        if suffix == '.txt':
            # Try UTF-8, then latin-1
            for encoding in ('utf-8', 'latin-1'):
                try:
                    with f.open('r', encoding=encoding, errors='ignore') as fp:
                        content = fp.read(max_chars)
                        return _clean_text(content)
                except UnicodeDecodeError:
                    continue
                except Exception as e:
                    return f"[Error reading file: {str(e)}]"
            return f"[Unreadable file: {suffix}]"

        elif suffix == '.pdf':
            # Try text extraction with pdfplumber, fallback to PyPDF2, then OCR
            text = ""
            if pdfplumber:
                try:
                    with pdfplumber.open(str(f)) as pdf:
                        for page in pdf.pages:
                            text += page.extract_text() or ""
                            if len(text) >= max_chars:
                                break
                except Exception:
                    text = ""
            if not text and PyPDF2:
                try:
                    with open(f, 'rb') as fp:
                        reader = PyPDF2.PdfReader(fp)
                        for page in reader.pages:
                            text += page.extract_text() or ""
                            if len(text) >= max_chars:
                                break
                except Exception:
                    text = ""
            if not text and pytesseract and Image:
                # OCR fallback: render each page as image and OCR
                try:
                    import pdf2image
                    images = pdf2image.convert_from_path(str(f))
                    for img in images:
                        text += pytesseract.image_to_string(img)
                        if len(text) >= max_chars:
                            break
                except Exception:
                    text = ""
            if text:
                return _clean_text(text)[:max_chars]
            else:
                return "[Could not extract text from PDF]"

        elif suffix in ('.docx',):
            if docx:
                try:
                    doc = docx.Document(str(f))
                    text = "\n".join(p.text for p in doc.paragraphs)
                    return _clean_text(text)[:max_chars]
                except Exception as e:
                    return f"[Error reading DOCX: {str(e)}]"
            else:
                return "[python-docx not installed]"

        elif suffix in ('.doc',):
            # Use antiword via subprocess for legacy .doc files
            try:
                result = subprocess.run([
                    'antiword',
                    str(f)
                ],
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    check=False)
                text = result.stdout.decode('utf-8', errors='ignore')
                if text:
                    return _clean_text(text)[:max_chars]
                return "[No text extracted from DOC]"
            except FileNotFoundError:
                return "[antiword not installed]"
            except Exception as e:
                return f"[Error reading DOC: {str(e)}]"

        elif suffix in ('.pptx',):
            if pptx:
                try:
                    prs = pptx.Presentation(str(f))
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text.append(shape.text)
                    return _clean_text("\n".join(text))[:max_chars]
                except Exception as e:
                    return f"[Error reading PPTX: {str(e)}]"
            else:
                return "[python-pptx not installed]"

        elif suffix in ('.xlsx',):
            if openpyxl:
                try:
                    wb = openpyxl.load_workbook(str(f), read_only=True, data_only=True)
                    text = []
                    for ws in wb.worksheets:
                        for row in ws.iter_rows(values_only=True):
                            for cell in row:
                                if cell is not None:
                                    text.append(str(cell))
                                    if sum(len(t) for t in text) >= max_chars:
                                        break
                            if sum(len(t) for t in text) >= max_chars:
                                break
                        if sum(len(t) for t in text) >= max_chars:
                            break
                    return _clean_text(" ".join(text))[:max_chars]
                except Exception as e:
                    return f"[Error reading XLSX: {str(e)}]"
            else:
                return "[openpyxl not installed]"

        else:
            return f"[Unsupported file type: {suffix}]"
    except Exception as e:
        return f"[Error extracting content: {str(e)}]"
 

def _clean_text(text: str) -> str:
    """
    Clean extracted text: collapse whitespace, strip control chars, etc.
    """
    import re
    text = re.sub(r'[\r\n]+', '\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    text = text.strip()
    return text
 

 main
